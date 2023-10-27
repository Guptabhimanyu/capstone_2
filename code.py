# import streamlit as st
# from dotenv import load_dotenv
# from PyPDF2 import PdfReader
# from docx import Document
# from openpyxl import load_workbook  # Add this import for Excel
# from pptx import Presentation  # Add this import for presentations
# from langchain.text_splitter import CharacterTextSplitter
# from langchain.embeddings.huggingface import HuggingFaceEmbeddings
# from langchain.vectorstores import FAISS
# from langchain.chains.question_answering import load_qa_chain
# from langchain import HuggingFaceHub


# # Function to extract text from Excel (XLSX) file


# def extract_text_from_excel(xlsx_file):
#     workbook = load_workbook(xlsx_file)
#     text = ""
#     for sheet in workbook.sheetnames:
#         for row in workbook[sheet].iter_rows(values_only=True):
#             text += " ".join(map(str, row)) + "\n"
#     return text

# # Function to extract text from PowerPoint (PPTX) file


# def extract_text_from_presentation(pptx_file):
#     prs = Presentation(pptx_file)
#     text = ""
#     for slide in prs.slides:
#         for shape in slide.shapes:
#             if shape.has_text_frame:
#                 text += shape.text + "\n"
#     return text


# def main():
#     load_dotenv()
#     st.set_page_config(page_title="Ask your Document")
#     st.header("Ask Your Document")

#     # Sidebar
#     uploaded_file = st.sidebar.file_uploader(
#         "Upload your document (PDF, Word, Excel, or Presentation)", type=["pdf", "docx", "xlsx", "pptx"])
#     user_question = st.sidebar.text_input(
#         "Ask a question about your document:")

#     # Main panel
#     if uploaded_file is not None:
#         if uploaded_file.type == "application/pdf":
#             # Handle PDF
#             pdf_reader = PdfReader(uploaded_file)
#             text = ""
#             for page in pdf_reader.pages:
#                 text += page.extract_text()
#         elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
#             # Handle Word document
#             text = extract_text_from_word_docx(uploaded_file)
#         elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
#             # Handle Excel file
#             text = extract_text_from_excel(uploaded_file)
#         elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
#             # Handle Presentation (PPTX) file
#             text = extract_text_from_presentation(uploaded_file)
#         else:
#             st.error(
#                 "Unsupported file format. Please upload a PDF, Word document, Excel file, or Presentation.")
#             return

#         # Split text into chunks and create embeddings
#         text_splitter = CharacterTextSplitter(
#             separator="\n",
#             chunk_size=1000,
#             chunk_overlap=200,
#             length_function=len
#         )
#         chunks = text_splitter.split_text(text)

#         embeddings = HuggingFaceEmbeddings()
#         knowledge_base = FAISS.from_texts(chunks, embeddings)

#         if user_question:
#             docs = knowledge_base.similarity_search(user_question)
#             llm = HuggingFaceHub(repo_id="google/flan-t5-large",
#                                  model_kwargs={"temperature": 5, "max_length": 10000})
#             chain = load_qa_chain(llm, chain_type="stuff")
#             response = chain.run(input_documents=docs, question=user_question)

#             st.write(response)


# if __name__ == '__main__':
#     main()

import streamlit as st
from dotenv import load_dotenv
from PyPDF2 import PdfReader
from docx import Document
from openpyxl import load_workbook  # Add this import for Excel
from pptx import Presentation  # Add this import for presentations
from langchain.text_splitter import CharacterTextSplitter
from langchain.embeddings.huggingface import HuggingFaceEmbeddings
from langchain.vectorstores import FAISS
from langchain.chains.question_answering import load_qa_chain
from langchain import HuggingFaceHub

# Function to extract text from Excel (XLSX) file


def extract_text_from_excel(xlsx_file):
    workbook = load_workbook(xlsx_file)
    text = ""
    for sheet in workbook.sheetnames:
        for row in workbook[sheet].iter_rows(values_only=True):
            text += " ".join(map(str, row)) + "\n"
    return text

# Function to extract text from PowerPoint (PPTX) file


def extract_text_from_presentation(pptx_file):
    prs = Presentation(pptx_file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text += shape.text + "\n"
    return text


def main():
    load_dotenv()
    st.set_page_config(page_title="Ask your Document")
    st.header("Ask Your Document")

    # Sidebar
    uploaded_file = st.sidebar.file_uploader(
        "Upload your document (PDF, Word, Excel, or Presentation)", type=["pdf", "docx", "xlsx", "pptx"])
    user_question = st.sidebar.text_input(
        "Ask a question about your document:")

    # Main panel
    if uploaded_file is not None:
        if uploaded_file.type == "application/pdf":
            # Handle PDF
            pdf_reader = PdfReader(uploaded_file)
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text()
        elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            # Handle Word document
            text = extract_text_from_word_docx(uploaded_file)
        elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
            # Handle Excel file
            text = extract_text_from_excel(uploaded_file)
        elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            # Handle Presentation (PPTX) file
            text = extract_text_from_presentation(uploaded_file)
        else:
            st.error(
                "Unsupported file format. Please upload a PDF, Word document, Excel file, or Presentation.")
            return

        # Split text into chunks and create embeddings
        text_splitter = CharacterTextSplitter(
            separator="\n",
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len
        )
        chunks = text_splitter.split_text(text)

        embeddings = HuggingFaceEmbeddings()
        knowledge_base = FAISS.from_texts(chunks, embeddings)

        if user_question:
            docs = knowledge_base.similarity_search(user_question)
            llm = HuggingFaceHub(repo_id="google/flan-t5-large",
                                 model_kwargs={"temperature": 5, "max_length": 10000})
            # Change the chain type to one of the supported values
            chain = load_qa_chain(llm, chain_type="stuff")
            response = chain.run(input_documents=docs, question=user_question)

            st.write(response)


if __name__ == '__main__':
    main()
