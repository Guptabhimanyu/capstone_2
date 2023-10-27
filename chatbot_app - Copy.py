import streamlit as st
from dotenv import load_dotenv
from PyPDF2 import PdfReader
from docx import Document
from openpyxl import load_workbook  # Add this import for Excel
from pptx import Presentation
import os
import base64
import time
from transformers import AutoTokenizer, AutoModelForSeq2SeqLM
from transformers import pipeline
import torch
import textwrap
from langchain.document_loaders import PyPDFLoader, DirectoryLoader, PDFMinerLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.embeddings import SentenceTransformerEmbeddings
from langchain.vectorstores import Chroma
from langchain.llms import HuggingFacePipeline
from langchain.chains import RetrievalQA
from constants import CHROMA_SETTINGS
from streamlit_chat import message
from langchain.text_splitter import CharacterTextSplitter
from langchain.embeddings.huggingface import HuggingFaceEmbeddings
from langchain.vectorstores import FAISS
from langchain.chains.question_answering import load_qa_chain
from langchain import HuggingFaceHub

import torch
from transformers import AutoTokenizer, AutoModelForSeq2SeqLM

# Set the device to CPU
device = torch.device('cpu')

checkpoint = "MBZUAI/LaMini-T5-738M"
print(f"Checkpoint path: {checkpoint}")

tokenizer = AutoTokenizer.from_pretrained(checkpoint)
base_model = AutoModelForSeq2SeqLM.from_pretrained(
    checkpoint, torch_dtype=torch.float32)
base_model.to(device)

# # tokenizer = AutoTokenizer.from_pretrained(checkpoint)
# # base_model = AutoModelForSeq2SeqLM.from_pretrained(
# #     checkpoint,
# #     device=device,  # Use the 'device' argument to specify the device
# #     torch_dtype=torch.float32
# # )


# persist_directory = "db"

# @st.cache_resource
# def data_ingestion():
#     for root, dirs, files in os.walk("docs"):
#         for file in files:
#             if file.endswith(".pdf"):
#                 print(file)
#                 loader = PDFMinerLoader(os.path.join(root, file))
#     documents = loader.load()
#     text_splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=500)
#     texts = text_splitter.split_documents(documents)
#     #create embeddings here
#     embeddings = SentenceTransformerEmbeddings(model_name="all-MiniLM-L6-v2")
#     #create vector store here
#     db = Chroma.from_documents(texts, embeddings, persist_directory=persist_directory, client_settings=CHROMA_SETTINGS)
#     db.persist()
#     db=None

# @st.cache_resource
# def llm_pipeline():
#     pipe = pipeline(
#         'text2text-generation',
#         model = base_model,
#         tokenizer = tokenizer,
#         max_length = 256,
#         do_sample = True,
#         temperature = 0.3,
#         top_p= 0.95,
#         device=device
#     )
#     local_llm = HuggingFacePipeline(pipeline=pipe)
#     return local_llm

# @st.cache_resource
# def qa_llm():
#     llm = llm_pipeline()
#     embeddings = SentenceTransformerEmbeddings(model_name="all-MiniLM-L6-v2")
#     db = Chroma(persist_directory="db", embedding_function = embeddings, client_settings=CHROMA_SETTINGS)
#     retriever = db.as_retriever()
#     qa = RetrievalQA.from_chain_type(
#         llm = llm,
#         chain_type = "stuff",
#         retriever = retriever,
#         return_source_documents=True
#     )
#     return qa

# def process_answer(instruction):
#     response = ''
#     instruction = instruction
#     qa = qa_llm()
#     generated_text = qa(instruction)
#     answer = generated_text['result']
#     return answer

# def get_file_size(file):
#     file.seek(0, os.SEEK_END)
#     file_size = file.tell()
#     file.seek(0)
#     return file_size

@st.cache_data
# function to display the PDF of a given file
def displayPDF(file):
    # Opening file from file path
    with open(file, "rb") as f:
        base64_pdf = base64.b64encode(f.read()).decode('utf-8')

    # Embedding PDF in HTML
    pdf_display = F'<iframe src="data:application/pdf;base64,{base64_pdf}" width="100%" height="600" type="application/pdf"></iframe>'

    # Displaying File
    st.markdown(pdf_display, unsafe_allow_html=True)

# Display conversation history using Streamlit messages


def display_conversation(history):
    for i in range(len(history["generated"])):
        message(history["past"][i], is_user=True, key=str(i) + "_user")
        message(history["generated"][i], key=str(i))


def main():
    st.markdown("<h1 style='text-align: center; color: blue;'>Chat with your PDF 🦜📄 </h1>",
                unsafe_allow_html=True)
    st.markdown("<h3 style='text-align: center; color: grey;'>Built by <a href='https://github.com/AIAnytime'>AI Anytime with ❤️ </a></h3>", unsafe_allow_html=True)

    st.markdown("<h2 style='text-align: center; color:red;'>Upload your PDF 👇</h2>",
                unsafe_allow_html=True)

    # uploaded_file = st.file_uploader("", type=["pdf"])


def extract_text_from_excel(xlsx_file):
    workbook = load_workbook(xlsx_file)
    text = ""
    for sheet in workbook.sheetnames:
        for row in workbook[sheet].iter_rows(values_only=True):
            text += " ".join(map(str, row)) + "\n"
    return text

# Function to extract text from PowerPoint (PPTX) file


def extract_text_from_presentation(pptx_file):
    prs = Presentation(pptx_file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text += shape.text + "\n"
    return text


def main():
    load_dotenv()
    st.set_page_config(page_title="Ask your Document")
    st.header("Ask Your Document")

    uploaded_file = st.file_uploader("Upload your document (PDF, Word, Excel, or Presentation)", type=[
                                     "pdf", "docx", "xlsx", "pptx"])

    # if uploaded_file is not None:
    #     file_details = {
    #         "Filename": uploaded_file.name,
    #         "File size": get_file_size(uploaded_file)
    #     }
    #     filepath = "docs/"+uploaded_file.name
    #     with open(filepath, "wb") as temp_file:
    #             temp_file.write(uploaded_file.read())

    #     col1, col2= st.columns([1,2])
    #     with col1:
    #         st.markdown("<h4 style color:black;'>File details</h4>", unsafe_allow_html=True)
    #         st.json(file_details)
    #         st.markdown("<h4 style color:black;'>File preview</h4>", unsafe_allow_html=True)
    #         pdf_view = displayPDF(filepath)

    #     with col2:
    #         with st.spinner('Embeddings are in process...'):
    #             ingested_data = data_ingestion()
    #         st.success('Embeddings are created successfully!')
    #         st.markdown("<h4 style color:black;'>Chat Here</h4>", unsafe_allow_html=True)

    if uploaded_file is not None:
        if uploaded_file.type == "application/pdf":
            # Handle PDF
            pdf_reader = PdfReader(uploaded_file)
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text()
        elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            # Handle Word document
            text = extract_text_from_word_docx(uploaded_file)
        elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
            # Handle Excel file
            text = extract_text_from_excel(uploaded_file)
        elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            # Handle Presentation (PPTX) file
            text = extract_text_from_presentation(uploaded_file)
        else:
            st.error(
                "Unsupported file format. Please upload a PDF, Word document, Excel file, or Presentation.")
            return

            # # Initialize session state for generated responses and past messages
            # if "generated" not in st.session_state:
            #     st.session_state["generated"] = ["I am ready to help you"]
            # if "past" not in st.session_state:
            #     st.session_state["past"] = ["Hey there!"]

        text_splitter = CharacterTextSplitter(
            separator="\n",
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len
        )
        chunks = text_splitter.split_text(text)
        embeddings = HuggingFaceEmbeddings()
        knowledge_base = FAISS.from_texts(chunks, embeddings)

        user_question = st.text_input(
            "Ask a question about your document:", key="input")

        if user_question:
            docs = knowledge_base.similarity_search(user_question)
            llm = HuggingFaceHub(repo_id="google/flan-t5-large",
                                 model_kwargs={"temperature": 5, "max_length": 10000})
            chain = load_qa_chain(llm, chain_type="stuff")
            response = chain.run(input_documents=docs,
                                 question=user_question)

            st.write(response)

            # Search the database for a response based on user input and update session state
            # if user_input:
            #     answer = process_answer({'query': user_input})
            #     st.session_state["past"].append(user_input)
            #     response = answer
            #     st.session_state["generated"].append(response)

            # # Display conversation history using Streamlit messages
            # if st.session_state["generated"]:
            #     display_conversation(st.session_state)


if __name__ == "__main__":
    main()
